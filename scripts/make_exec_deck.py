"""Generate a 2-slide executive deck (Architecture + Data Model) for Zava.

Run:
    python scripts/make_exec_deck.py
Produces: Zava_Exec_Overview.pptx in the repo root.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ---- Palette -------------------------------------------------------------
INK = RGBColor(0x1B, 0x1F, 0x2A)        # near-black text
SUBINK = RGBColor(0x5A, 0x61, 0x72)     # muted text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF4, 0xF6, 0xFB)         # slide background
CARD = RGBColor(0xFF, 0xFF, 0xFF)

BLUE = RGBColor(0x2F, 0x6F, 0xED)       # app / orchestration
PURPLE = RGBColor(0x6B, 0x4D, 0xE0)     # agents
TEAL = RGBColor(0x12, 0x9A, 0x8F)       # grounding/data
AMBER = RGBColor(0xE2, 0x8A, 0x00)      # action / highlight
SLATE = RGBColor(0x3C, 0x46, 0x5B)      # neutral box
GREEN = RGBColor(0x1E, 0x9E, 0x5A)
ROSE = RGBColor(0xD9, 0x3A, 0x5B)

EMU = 914400  # per inch
SW, SH = int(13.333 * EMU), int(7.5 * EMU)


def _in(v: float) -> int:
    return int(v * EMU)


def add_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide


def _no_shadow(shape):
    shape.shadow.inherit = False


def box(slide, x, y, w, h, fill=CARD, line=None, line_w=1.0, radius=0.10, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, _in(x), _in(y), _in(w), _in(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    _no_shadow(sp)
    # corner radius
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=2, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color)."""
    tb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, size, bold, color) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
    return tb


def label(slide, x, y, w, h, s, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER):
    text(slide, x, y, w, h, [[(s, size, bold, color)]], align=align, anchor=MSO_ANCHOR.MIDDLE)


def chip(slide, x, y, w, h, title, sub, fill, title_c=WHITE, sub_c=None):
    sp = box(slide, x, y, w, h, fill=fill, radius=0.16)
    sub_c = sub_c or RGBColor(0xEC, 0xF0, 0xFB)
    runs = [[(title, 12, True, title_c)]]
    if sub:
        runs.append([(sub, 8.5, False, sub_c)])
    text(slide, x, y, w, h, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         space_after=1, line_spacing=0.95)
    return sp


def arrow(slide, x1, y1, x2, y2, color=SUBINK, w=1.75):
    cn = slide.shapes.add_connector(2, _in(x1), _in(y1), _in(x2), _in(y2))  # straight
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    _no_shadow(cn)
    line = cn.line._get_or_add_ln()
    tail = line.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    line.append(tail)
    return cn


def title_bar(slide, kicker, title):
    box(slide, 0, 0, 13.333, 1.18, fill=INK, radius=0, shape=MSO_SHAPE.RECTANGLE)
    box(slide, 0, 1.12, 13.333, 0.06, fill=AMBER, radius=0, shape=MSO_SHAPE.RECTANGLE)
    text(slide, 0.55, 0.16, 12.3, 0.4, [[(kicker, 12, True, AMBER)]])
    text(slide, 0.55, 0.44, 12.3, 0.6, [[(title, 24, True, WHITE)]])


# =========================================================================
# SLIDE 1 — Use case + Architecture
# =========================================================================
def build_slide1(prs):
    s = add_slide(prs)
    title_bar(s, "ZAVA · INSIGHTS-TO-ACTION AGENT", "From a plain-English question to a prioritised business action")

    # --- Left rail: use case / what & why ---
    box(s, 0.45, 1.45, 3.5, 5.5, fill=CARD, line=RGBColor(0xE0, 0xE4, 0xEE), radius=0.05)
    text(s, 0.7, 1.62, 3.1, 0.4, [[("THE USE CASE", 12, True, BLUE)]])
    text(s, 0.7, 2.0, 3.05, 1.5, [
        [("Retail decisions are ", 11, False, INK), ("cross-domain", 11, True, INK),
         (". A good call needs three facts at once:", 11, False, INK)],
        [("•  Is it selling more or less?", 10.5, False, SUBINK)],
        [("•  Can we fulfil it from stock?", 10.5, False, SUBINK)],
        [("•  Is a campaign behind it, and is it working?", 10.5, False, SUBINK)],
    ], space_after=4, line_spacing=1.02)

    text(s, 0.7, 3.7, 3.05, 0.4, [[("WHAT WE ARE DOING", 12, True, PURPLE)]])
    text(s, 0.7, 4.05, 3.05, 1.6, [
        [("A team of AI agents reads ", 10.5, False, INK), ("live", 10.5, True, INK),
         (" Sales, Inventory & Marketing data, then an ", 10.5, False, INK),
         ("Action agent", 10.5, True, AMBER),
         (" turns the insights into a ranked set of actions.", 10.5, False, INK)],
    ], line_spacing=1.05)

    text(s, 0.7, 5.5, 3.05, 0.4, [[("WHAT WE ACHIEVE", 12, True, TEAL)]])
    text(s, 0.7, 5.85, 3.05, 1.05, [
        [("•  Hours of cross-team analysis  →  seconds", 10.5, False, INK)],
        [("•  Every number is grounded & auditable", 10.5, False, INK)],
        [("•  Each action has an owner + priority", 10.5, False, INK)],
    ], space_after=3, line_spacing=1.0)

    # --- Right: architecture flow ---
    ax = 4.25
    # column headers
    text(s, ax, 1.4, 8.6, 0.3, [[("HOW IT WORKS  ·  question flows left to right", 11, True, SUBINK)]])

    # User
    chip(s, ax, 3.5, 1.15, 1.0, "User", "asks in chat", SLATE)

    # Chat app (Container Apps)
    box(s, ax + 1.45, 1.95, 2.35, 4.05, fill=RGBColor(0xEAF0FE >> 16 & 255 and 0xEA, 0xF0, 0xFE),
        line=BLUE, line_w=1.25, radius=0.06)
    text(s, ax + 1.45, 2.02, 2.35, 0.3, [[("AZURE CONTAINER APPS", 8.5, True, BLUE)]], align=PP_ALIGN.CENTER)
    chip(s, ax + 1.62, 2.35, 2.0, 0.85, "Chat app", "FastAPI · live UI + charts", BLUE)
    chip(s, ax + 1.62, 3.35, 2.0, 0.95, "Orchestrator", "Agent Framework (Magentic)", BLUE)
    chip(s, ax + 1.62, 4.45, 2.0, 0.85, "Manager", "plans · gpt-4.1-mini", RGBColor(0x1D, 0x4E, 0xD8))

    # Foundry agents (now incl. Intent Detector + Response Generator)
    box(s, ax + 4.05, 1.5, 2.4, 5.36, fill=RGBColor(0xF0, 0xEC, 0xFB), line=PURPLE, line_w=1.25, radius=0.05)
    text(s, ax + 4.05, 1.55, 2.4, 0.3, [[("MICROSOFT FOUNDRY · HOSTED AGENTS", 8.0, True, PURPLE)]], align=PP_ALIGN.CENTER)
    chip(s, ax + 4.22, 1.9, 2.05, 0.62, "Intent Detector", "chat vs business", SLATE)
    chip(s, ax + 4.22, 2.56, 2.05, 0.54, "Sales agent", "trends & revenue", PURPLE)
    chip(s, ax + 4.22, 3.14, 2.05, 0.54, "Inventory agent", "stock & reorder", PURPLE)
    chip(s, ax + 4.22, 3.72, 2.05, 0.54, "Marketing agent", "campaigns & ROI", PURPLE)
    chip(s, ax + 4.22, 4.32, 2.05, 0.72, "Action agent", "actions + chart\n(Code Interpreter)", AMBER)
    chip(s, ax + 4.22, 5.1, 2.05, 0.7, "Response Gen", "writes the final reply", GREEN)
    text(s, ax + 4.22, 5.86, 2.05, 0.92, [
        [("Greetings skip straight to Response Gen; business questions fan out, then return here.", 7.6, False, SUBINK)],
    ], align=PP_ALIGN.CENTER, line_spacing=0.92)

    # Data layer
    box(s, ax + 6.7, 1.95, 1.9, 4.05, fill=RGBColor(0xE6, 0xF5, 0xF3), line=TEAL, line_w=1.25, radius=0.06)
    text(s, ax + 6.7, 2.02, 1.9, 0.45, [[("GROUNDING", 8.5, True, TEAL)]], align=PP_ALIGN.CENTER)
    chip(s, ax + 6.85, 2.45, 1.6, 0.7, "MCP tools", "typed, read-only", TEAL)
    chip(s, ax + 6.85, 3.3, 1.6, 0.75, "Cosmos DB", "sales · inv · mktg", TEAL)
    chip(s, ax + 6.85, 4.2, 1.6, 0.75, "Foundry IQ", "campaign briefs", GREEN)
    chip(s, ax + 6.85, 5.1, 1.6, 0.72, "Code sandbox", "consolidated math", SLATE)

    # arrows
    arrow(s, ax + 1.15, 4.0, ax + 1.45, 4.0, color=SLATE)
    arrow(s, ax + 3.62, 4.0, ax + 4.05, 4.0, color=BLUE)
    arrow(s, ax + 6.27, 3.6, ax + 6.7, 3.6, color=PURPLE)
    # action -> final note handled by text

    # bottom strip: example
    box(s, 0.45, 7.02, 12.45, 0.42, fill=INK, radius=0.5)
    text(s, 0.6, 7.04, 12.2, 0.4,
         [[("Example:  ", 10.5, True, AMBER),
           ("“Garden sales are trending up — do we have stock and a live campaign to support it? What should we do?”  →  one grounded answer + ranked actions.",
            10.5, False, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


# =========================================================================
# SLIDE 2 — Data model (plain-language, customer-friendly)
# =========================================================================
def build_slide2(prs):
    s = add_slide(prs)
    title_bar(s, "ZAVA · THE DATA BEHIND THE ANSWERS", "Three everyday business records the agents read")

    text(s, 0.55, 1.3, 12.2, 0.45,
         [[("The assistant never guesses. ", 12, True, INK),
           ("Every answer is built from three simple business records below — the same data your store, supply and marketing teams already use.",
            12, False, INK)]])

    cards = [
        dict(
            x=0.45, color=BLUE, icon="🧾", name="SALES  ·  what we sold",
            plain="One line for every sale we ring up, across all stores and online.",
            example=("Example:  On 12 May, the Seattle store sold "
                     "4 cans of Premium Interior Paint to a DIY customer "
                     "for $180 (we kept $70 profit)."),
            captures=["Which product, store and region",
                      "How many units and the revenue",
                      "The profit we made on it",
                      "Customer type: DIY, Pro or Contractor",
                      "12 months of history, month by month"],
            answers="“What are our best sellers? Is paint growing or cooling? Which stores are ahead?”"),
        dict(
            x=4.62, color=TEAL, icon="📦", name="INVENTORY  ·  what we have",
            plain="How much stock sits in each warehouse, refreshed every week.",
            example=("Example:  This week the Seattle warehouse has only "
                     "3 cans of Premium Interior Paint left — below the safe level, "
                     "so it needs reordering now."),
            captures=["Stock on hand per warehouse",
                      "A health flag: healthy / low / overstock",
                      "The reorder point and suggested order qty",
                      "“Weeks of cover” — how long stock will last",
                      "4 suppliers feeding 6 warehouses"],
            answers="“What is running low and needs reordering? Where do we have too much?”"),
        dict(
            x=8.79, color=PURPLE, icon="📣", name="MARKETING  ·  what we promote",
            plain="Every promotional campaign — what is live, planned, and how past ones did.",
            example=("Example:  ‘Spring Paint Sale’ is live in Seattle, 15% off paint, "
                     "$40k budget. Last year’s version returned $1.42 for every $1 spent."),
            captures=["Campaign status: live / planned / finished",
                      "Which products, stores and discount",
                      "Budget, spend, views and clicks",
                      "Return on investment (on past campaigns)",
                      "Written brief + lessons learned (the “why”)"],
            answers="“Which promos are running? Did they pay off? What should we change?”"),
    ]
    cw = 4.0
    for c in cards:
        x = c["x"]
        box(s, x, 1.9, cw, 4.55, fill=CARD, line=RGBColor(0xE0, 0xE4, 0xEE), radius=0.04)
        # header band
        box(s, x, 1.9, cw, 0.6, fill=c["color"], radius=0.14)
        text(s, x + 0.2, 1.9, cw - 0.4, 0.6, [[(c["icon"] + "   " + c["name"], 12.5, True, WHITE)]],
             anchor=MSO_ANCHOR.MIDDLE)
        # plain definition
        text(s, x + 0.2, 2.6, cw - 0.4, 0.55, [[(c["plain"], 10.5, True, INK)]], line_spacing=1.0)
        # example callout
        box(s, x + 0.16, 3.18, cw - 0.32, 0.92, fill=RGBColor(0xF3, 0xF6, 0xFD), radius=0.08)
        text(s, x + 0.3, 3.22, cw - 0.6, 0.86, [[(c["example"], 9.2, False, RGBColor(0x33, 0x3B, 0x52))]],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.98)
        # what it captures
        runs = [[("WHAT IT TELLS US", 8.8, True, c["color"])]]
        for f in c["captures"]:
            runs.append([("•  " + f, 9.3, False, INK)])
        text(s, x + 0.2, 4.18, cw - 0.4, 1.6, runs, space_after=2, line_spacing=1.0)
        # answers band
        box(s, x + 0.16, 5.82, cw - 0.32, 0.55, fill=RGBColor(0xFB, 0xF3, 0xE2), radius=0.12)
        text(s, x + 0.28, 5.84, cw - 0.55, 0.52,
             [[("Answers:  ", 8.8, True, RGBColor(0x9A, 0x5E, 0x00)),
               (c["answers"], 8.6, False, RGBColor(0x6B, 0x45, 0x10))]],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.92)

    # bottom: how they connect -> the value
    box(s, 0.45, 6.6, 12.45, 0.78, fill=INK, radius=0.1)
    text(s, 0.65, 6.62, 12.1, 0.74, [
        [("The magic is joining all three:  ", 11, True, AMBER),
         ("paint sales are cooling, Seattle is nearly out of paint, and a paint promo is live — "
          "so the assistant flags it and recommends an emergency restock before we promote a product we can’t sell.",
          10.5, False, WHITE)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    return s


# =========================================================================
# SLIDE 3 — Where the data is stored (Cosmos DB vs Foundry IQ)
# =========================================================================
def build_slide3(prs):
    s = add_slide(prs)
    title_bar(s, "ZAVA · WHERE THE DATA LIVES", "Numbers in a database, written knowledge in Foundry IQ")

    text(s, 0.55, 1.32, 12.2, 0.5,
         [[("We keep data in the place that fits its shape. ", 12, True, INK),
           ("Exact numbers go in a database for precise math; written documents go in Foundry IQ so the agent can search them by meaning.",
            12, False, INK)]])

    # ---- LEFT: Cosmos DB (numbers) ----
    lx = 0.45
    box(s, lx, 2.0, 6.05, 4.55, fill=CARD, line=TEAL, line_w=1.5, radius=0.04)
    box(s, lx, 2.0, 6.05, 0.7, fill=TEAL, radius=0.1)
    text(s, lx + 0.25, 2.0, 5.6, 0.7,
         [[("🗄   AZURE COSMOS DB", 14, True, WHITE), ("   ·  the numbers", 11, False, RGBColor(0xD7, 0xF0, 0xEC))]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, lx + 0.25, 2.78, 5.55, 0.4,
         [[("Structured records — rows of exact figures the agents add up and rank.", 10, True, INK)]])

    cos = [
        ("Sales records", "every sale: product, store, units, revenue, profit", BLUE),
        ("Inventory snapshots", "weekly stock per warehouse + health & reorder level", TEAL),
        ("Marketing campaigns", "the KPIs: budget, spend, views, clicks, ROI, dates", PURPLE),
    ]
    yy = 3.3
    for name, desc, col in cos:
        box(s, lx + 0.25, yy, 5.55, 0.92, fill=RGBColor(0xF3, 0xF7, 0xFB), radius=0.08)
        box(s, lx + 0.25, yy, 0.14, 0.92, fill=col, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
        text(s, lx + 0.55, yy + 0.06, 5.2, 0.85, [
            [(name, 11, True, INK)],
            [(desc, 9.3, False, SUBINK)],
        ], space_after=1, line_spacing=0.95, anchor=MSO_ANCHOR.MIDDLE)
        yy += 1.02

    # ---- RIGHT: Foundry IQ (documents) ----
    rx = 6.85
    box(s, rx, 2.0, 6.05, 4.55, fill=CARD, line=GREEN, line_w=1.5, radius=0.04)
    box(s, rx, 2.0, 6.05, 0.7, fill=GREEN, radius=0.1)
    text(s, rx + 0.25, 2.0, 5.6, 0.7,
         [[("📚   FOUNDRY IQ", 14, True, WHITE), ("   ·  the written knowledge", 11, False, RGBColor(0xD8, 0xF1, 0xE4))]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, rx + 0.25, 2.78, 5.55, 0.4,
         [[("Documents the agent searches ", 10, True, INK), ("by meaning", 10, True, GREEN),
           (" — the “why” behind the numbers.", 10, True, INK)]])

    iq = [
        ("Campaign briefs", "the plan: goal, target audience, offer, creative direction", GREEN),
        ("Post-mortems", "lessons learned: what worked, what to change next time", AMBER),
        ("Marketing playbooks", "best-practice guidance the team follows", SLATE),
    ]
    yy = 3.3
    for name, desc, col in iq:
        box(s, rx + 0.25, yy, 5.55, 0.92, fill=RGBColor(0xF1, 0xF8, 0xF4), radius=0.08)
        box(s, rx + 0.25, yy, 0.14, 0.92, fill=col, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
        text(s, rx + 0.55, yy + 0.06, 5.2, 0.85, [
            [(name, 11, True, INK)],
            [(desc, 9.3, False, SUBINK)],
        ], space_after=1, line_spacing=0.95, anchor=MSO_ANCHOR.MIDDLE)
        yy += 1.02

    # bottom strip
    box(s, 0.45, 6.68, 12.45, 0.72, fill=INK, radius=0.1)
    text(s, 0.65, 6.7, 12.1, 0.68, [
        [("Why two stores?  ", 11, True, AMBER),
         ("A database gives exact, auditable numbers; Foundry IQ understands plain-English documents. "
          "The Marketing agent uses BOTH — the live KPIs and the written lessons — to answer “did it work, and what should we change?”",
          10, False, WHITE)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    return s


# =========================================================================
# SLIDE 4 — The scenario, step by step
# =========================================================================
def build_slide4(prs):
    s = add_slide(prs)
    title_bar(s, "ZAVA · ONE QUESTION, START TO FINISH", "How the agents turn a question into a decision")

    # the question
    box(s, 0.45, 1.36, 12.45, 0.6, fill=RGBColor(0xEA, 0xF0, 0xFE), line=BLUE, line_w=1.25, radius=0.12)
    text(s, 0.7, 1.36, 12.0, 0.6, [
        [("The user asks:   ", 11, True, BLUE),
         ("“Garden — sorry, paint sales are slipping in Seattle. Do we have stock and a campaign in play? What should we do?”",
          12, True, INK)]],
        anchor=MSO_ANCHOR.MIDDLE)

    # intent routing band (conditional edge)
    box(s, 0.45, 2.04, 12.45, 0.46, fill=RGBColor(0xE7, 0xEA, 0xF1), line=SLATE, line_w=1.0, radius=0.12)
    text(s, 0.7, 2.04, 12.0, 0.46, [
        [("Intent Detector:  ", 10.5, True, SLATE),
         ("first it classifies the message — a greeting would get a quick friendly reply; this is a business question, so the Manager engages the specialists.",
          10, False, INK)]],
        anchor=MSO_ANCHOR.MIDDLE)

    steps = [
        ("1", "SALES AGENT", BLUE, "Looks at sales",
         "Paint revenue in Seattle has been sliding for 3 months — the category is cooling."),
        ("2", "INVENTORY AGENT", TEAL, "Checks stock",
         "Seattle warehouse has just 3 cans of the top paint left — below the safe reorder level."),
        ("3", "MARKETING AGENT", PURPLE, "Reviews campaigns",
         "A ‘Spring Paint Sale’ is LIVE on that exact paint — and last year’s post-mortem (Foundry IQ) warned stock ran out early."),
        ("4", "ACTION AGENT", AMBER, "Decides the moves",
         "Joins it all up into a ranked set of actions — each with an owner and priority — plus a chart of the numbers."),
    ]
    cw, gap, x0, y0, ch = 2.95, 0.22, 0.45, 2.58, 2.2
    for i, (n, name, col, role, body) in enumerate(steps):
        x = x0 + i * (cw + gap)
        box(s, x, y0, cw, ch, fill=CARD, line=RGBColor(0xE0, 0xE4, 0xEE), radius=0.05)
        box(s, x, y0, cw, 0.62, fill=col, radius=0.12)
        text(s, x + 0.15, y0, cw - 0.3, 0.62, [[(n + "   " + name, 11.5, True, WHITE)]],
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + 0.2, y0 + 0.72, cw - 0.4, 0.35, [[(role, 10, True, col)]])
        text(s, x + 0.2, y0 + 1.08, cw - 0.4, 1.05, [[(body, 9.7, False, INK)]], line_spacing=1.02)
        if i < 3:
            arrow(s, x + cw + 0.01, y0 + ch / 2, x + cw + gap - 0.01, y0 + ch / 2, color=SUBINK, w=2.0)

    # the action result band
    box(s, 0.45, 5.0, 12.45, 1.5, fill=RGBColor(0xFB, 0xF3, 0xE2), line=AMBER, line_w=1.5, radius=0.06)
    text(s, 0.7, 5.06, 12.0, 0.4, [[("✓  RECOMMENDED ACTIONS  (what the leader sees)", 12, True, RGBColor(0x9A, 0x5E, 0x00))]])
    text(s, 0.7, 5.46, 12.0, 1.0, [
        [("1.  Emergency restock / transfer paint into Seattle this week   ", 10.5, True, INK),
         ("— Owner: Distributor Ops · Priority: High", 10, False, SUBINK)],
        [("2.  Hold or re-time the Spring Paint promo until stock lands   ", 10.5, True, INK),
         ("— Owner: Marketing · Priority: High", 10, False, SUBINK)],
        [("3.  Watch the paint downtrend; consider a margin review   ", 10.5, True, INK),
         ("— Owner: Merchandising · Priority: Medium", 10, False, SUBINK)],
    ], space_after=3, line_spacing=0.98)

    text(s, 0.45, 6.6, 12.45, 0.5,
         [[("The ", 10.5, True, INK), ("Response Generator", 10.5, True, GREEN),
           (" writes the final reply — keeping the Action agent’s chart — and every recommendation cites the exact numbers it came from.",
            10.5, False, INK)]],
         align=PP_ALIGN.CENTER)
    return s


def main():
    prs = Presentation()
    prs.slide_width = Emu(SW)
    prs.slide_height = Emu(SH)
    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)
    build_slide4(prs)
    out = "Zava_Exec_Overview.pptx"
    prs.save(out)
    print("Saved", out)


if __name__ == "__main__":
    main()
