"""Append two reference-architecture slides to Zava_Exec_Overview.pptx.

Slide A — Agentic AI reference architecture mapped to the Microsoft stack
          (the same logical layers as the shared "Agentic AI – Logical
          Architecture" diagram, populated with Microsoft products).
Slide B — How this Zava demo maps onto those same reference layers.

The design language (colours, fonts, header bar, tinted rounded cards) is
copied from the existing slides so the new ones look native to the deck.

Run:
    python scripts/add_reference_arch_slides.py
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DECK = "Zava_Exec_Overview.pptx"

# ---- palette (from the existing deck) ------------------------------------ #
BG        = RGBColor(0xF4, 0xF6, 0xFB)
DARK      = RGBColor(0x1B, 0x1F, 0x2A)
SLATE     = RGBColor(0x3C, 0x46, 0x5B)
MUTE      = RGBColor(0x5A, 0x61, 0x72)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x1B, 0x1F, 0x2A)
ORANGE    = RGBColor(0xE2, 0x8A, 0x00)
BLUE      = RGBColor(0x2F, 0x6F, 0xED)
PURPLE    = RGBColor(0x6B, 0x4D, 0xE0)
TEAL      = RGBColor(0x12, 0x9A, 0x8F)
GREEN     = RGBColor(0x1E, 0x9E, 0x5A)
CHIPSUB   = RGBColor(0xEC, 0xF0, 0xFB)

# light tints used for the band backgrounds
TINT = {
    "blue":   RGBColor(0xEA, 0xF0, 0xFE),
    "slate":  RGBColor(0xE7, 0xEA, 0xF1),
    "purple": RGBColor(0xF0, 0xEC, 0xFB),
    "teal":   RGBColor(0xE6, 0xF5, 0xF3),
    "green":  RGBColor(0xE5, 0xF5, 0xEB),
    "orange": RGBColor(0xFD, 0xF1, 0xDF),
}
SOLID = {
    "blue": BLUE, "slate": SLATE, "purple": PURPLE,
    "teal": TEAL, "green": GREEN, "orange": ORANGE,
}

FONT = "Segoe UI"


# --------------------------------------------------------------------------- #
# low-level helpers
# --------------------------------------------------------------------------- #
def _no_line(shape):
    shape.line.fill.background()


def round_rect(slide, x, y, w, h, fill, radius=0.08, line=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        _no_line(shp)
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    # adjust corner radius
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    return shp


def rect(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    _no_line(shp)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            wrap=True, space=None):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if space is not None:
            p.space_after = Pt(space)
        for (text, size, bold, color) in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tb


def header(slide, kicker, title):
    rect(slide, 0, 0, 13.333, 7.5, BG)              # base
    rect(slide, 0, 0, 13.333, 1.18, DARK)           # header bar
    rect(slide, 0, 1.12, 13.333, 0.06, ORANGE)      # accent line
    textbox(slide, 0.55, 0.16, 12.3, 0.4,
            [[(kicker, 12, True, ORANGE)]])
    textbox(slide, 0.55, 0.44, 12.3, 0.62,
            [[(title, 23, True, WHITE)]])


def footer(slide, bold_lead, rest):
    round_rect(slide, 0.45, 7.02, 12.45, 0.42, DARK, radius=0.25)
    textbox(slide, 0.6, 7.04, 12.2, 0.38,
            [[(bold_lead, 10.5, True, ORANGE), (rest, 10.5, False, WHITE)]],
            anchor=MSO_ANCHOR.MIDDLE)


def chip(slide, x, y, w, h, color, title, sub=None):
    round_rect(slide, x, y, w, h, color, radius=0.16)
    paras = [[(title, 10.5, True, WHITE)]]
    if sub:
        paras.append([(sub, 8, False, CHIPSUB)])
    textbox(slide, x + 0.08, y, w - 0.16, h, paras,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def band(slide, y, h, color_key, label, caption, chips):
    """One horizontal layer band: solid label box on the left + equal chips."""
    x = 0.45
    total_w = 12.45
    label_w = 2.65
    tint = TINT[color_key]
    solid = SOLID[color_key]
    # band background (light tint)
    round_rect(slide, x, y, total_w, h, tint, radius=0.10)
    # solid label box
    round_rect(slide, x, y, label_w, h, solid, radius=0.10)
    textbox(slide, x + 0.16, y, label_w - 0.3, h,
            [[(label, 11.5, True, WHITE)], [(caption, 8, False, CHIPSUB)]],
            anchor=MSO_ANCHOR.MIDDLE)
    # chips area
    cx0 = x + label_w + 0.18
    cw_total = (x + total_w) - cx0 - 0.18
    n = len(chips)
    gap = 0.14
    cw = (cw_total - gap * (n - 1)) / n
    ch = h - 0.22
    cy = y + 0.11
    cxs = cx0
    for (t, s) in chips:
        chip(slide, cxs, cy, cw, ch, solid, t, s)
        cxs += cw + gap


def add_layer_slide(prs, blank, kicker, title, bands, foot_lead, foot_rest):
    slide = prs.slides.add_slide(blank)
    header(slide, kicker, title)
    y = 1.40
    for (h, key, label, cap, chips) in bands:
        band(slide, y, h, key, label, cap, chips)
        y += h + 0.10
    footer(slide, foot_lead, foot_rest)
    return slide


# --------------------------------------------------------------------------- #
# build
# --------------------------------------------------------------------------- #
def _trim_to(prs, keep):
    """Drop any slides past the first ``keep`` so the script is re-runnable."""
    sld_id_lst = prs.slides._sldIdLst
    ids = list(sld_id_lst)
    for sld_id in ids[keep:]:
        rId = sld_id.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sld_id_lst.remove(sld_id)


def main():
    prs = Presentation(DECK)
    _trim_to(prs, 4)  # original deck has 4 slides; remove any prior re-runs
    blank = next(l for l in prs.slide_layouts if l.name == "Blank")

    # ---------------- Slide A : Microsoft-stack reference architecture ----- #
    bands_a = [
        (0.78, "blue", "EXPERIENCE", "how people & apps ask", [
            ("Microsoft 365 Copilot", "Teams · Outlook · Word"),
            ("Copilot Studio", "low-code agents"),
            ("Custom apps", "Web / FastAPI · API endpoints"),
            ("Agent 365 / Agent Store", "discover & share"),
        ]),
        (0.78, "slate", "TRUST & POLICY", "identity · gateway · safety", [
            ("Microsoft Entra Agent ID", "agent identity & AuthZ"),
            ("Azure API Management", "AI gateway · throttling"),
            ("Microsoft Purview", "DLP & governance"),
            ("Azure AI Content Safety", "+ Defender · responsible AI"),
        ]),
        (0.86, "purple", "ORCHESTRATION", "plan · run · remember", [
            ("Microsoft Foundry", "Agent Service · hosted agents"),
            ("Agent Framework", "Magentic multi-agent · A2A"),
            ("Azure Durable Functions", "durable workflow / state"),
            ("Memory", "short & long-term · threads"),
        ]),
        (0.78, "teal", "TOOLS & MCP", "discover · authorize · call", [
            ("API Center", "MCP registry / catalog"),
            ("APIM MCP gateway", "govern internal & external MCP"),
            ("Foundry MCP servers", "typed, read-only tools"),
            ("Logic Apps connectors", "1000+ enterprise apps"),
        ]),
        (0.78, "green", "KNOWLEDGE & DATA", "ground every answer", [
            ("Azure Cosmos DB", "operational records"),
            ("Foundry IQ / AI Search", "vector + semantic"),
            ("Microsoft Fabric", "OneLake · lakehouse"),
            ("Azure Databricks", "analytics / ML"),
        ]),
        (0.78, "orange", "MODELS", "reason & generate", [
            ("Foundry / Azure OpenAI", "GPT family LLMs"),
            ("Small language models", "Phi · cost-efficient"),
            ("Model router", "right model per task"),
            ("Code Interpreter", "consolidated math"),
        ]),
    ]
    add_layer_slide(
        prs, blank,
        "ZAVA · AGENTIC AI REFERENCE ARCHITECTURE",
        "The Microsoft agent platform behind this demo — at enterprise scale",
        bands_a,
        "Governance & observability span every layer:  ",
        "Azure Monitor · Application Insights · Foundry Evaluations · Content Safety · Cost Management (FinOps).",
    )

    # ---------------- Slide B : Zava demo mapped to the reference ---------- #
    bands_b = [
        (0.78, "blue", "EXPERIENCE", "how the user asks", [
            ("Chat app (FastAPI)", "live UI + charts"),
            ("Azure Container Apps", "serverless hosting"),
            ("Streaming UI", "plan · agents · trace"),
        ]),
        (0.78, "slate", "TRUST & POLICY", "identity & access", [
            ("Microsoft Entra ID", "DefaultAzureCredential"),
            ("Basic auth gate", "front-door protection"),
            ("Typed read-only MCP", "no write paths"),
        ]),
        (0.86, "purple", "ORCHESTRATION", "intent routing + Magentic + response", [
            ("Agent Framework", "Magentic multi-agent"),
            ("Intent Detector", "chat vs business · shared memory"),
            ("Manager · gpt-4.1-mini", "plans the specialists"),
            ("Sales·Inv·Mktg·Action", "Foundry hosted agents"),
            ("Response Generator", "final reply · keeps charts"),
        ]),
        (0.78, "teal", "TOOLS & MCP", "Cosmos-backed tools", [
            ("Products MCP", "sales & inventory"),
            ("Marketing MCP", "campaigns & ROI"),
            ("Code sandbox", "Code Interpreter"),
        ]),
        (0.78, "green", "KNOWLEDGE & DATA", "ground every number", [
            ("Azure Cosmos DB", "sales · inventory · marketing"),
            ("Foundry IQ", "briefs & post-mortems"),
            ("Bing web_search", "live context"),
        ]),
        (0.78, "orange", "MODELS", "reason & answer", [
            ("gpt-4.1-mini", "manager + agents"),
            ("Foundry agent models", "per-agent config"),
            ("Code Interpreter", "consolidated chart + math"),
        ]),
    ]
    add_layer_slide(
        prs, blank,
        "ZAVA · THIS DEMO ON THE REFERENCE ARCHITECTURE",
        "Every Zava component is a swap-in of the enterprise building block above it",
        bands_b,
        "Observability in the demo:  ",
        "Application Insights / OpenTelemetry tracing · Foundry Evaluations · red-team & guardrails modules.",
    )

    prs.save(DECK)
    print(f"Saved {DECK} — now {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
